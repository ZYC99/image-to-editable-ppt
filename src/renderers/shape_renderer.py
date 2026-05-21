from __future__ import annotations

from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.util import Pt

from src.ir.normalize import element_to_inches
from src.ir.schema import Element, Layout
from src.renderers.utils import hex_to_rgb, style_value


SHAPES = {
    "rect": MSO_SHAPE.RECTANGLE,
    "rectangle": MSO_SHAPE.RECTANGLE,
    "round_rect": MSO_SHAPE.ROUNDED_RECTANGLE,
    "rounded_rect": MSO_SHAPE.ROUNDED_RECTANGLE,
    "oval": MSO_SHAPE.OVAL,
    "triangle": MSO_SHAPE.ISOSCELES_TRIANGLE,
    "diamond": MSO_SHAPE.DIAMOND,
    "chevron": MSO_SHAPE.CHEVRON,
}


def render_shape(slide, layout: Layout, element: Element):
    left, top, width, height = element_to_inches(layout, element)
    shape_type = SHAPES.get(element.shape or style_value(element.style, "shape", "rect"), MSO_SHAPE.RECTANGLE)
    shape = slide.shapes.add_shape(shape_type, left, top, width, height)
    shape.name = element.id

    fill_color = style_value(element.style, "fill", "#FFFFFF")
    if fill_color == "transparent":
        shape.fill.background()
    else:
        shape.fill.solid()
        shape.fill.fore_color.rgb = hex_to_rgb(fill_color, "#FFFFFF")
        transparency = float(style_value(element.style, "fill_transparency", 0))
        shape.fill.transparency = max(0, min(100, transparency))

    line_color = style_value(element.style, "stroke", "transparent")
    if line_color == "transparent":
        shape.line.fill.background()
    else:
        shape.line.color.rgb = hex_to_rgb(line_color, "#000000")
        shape.line.width = Pt(float(style_value(element.style, "stroke_width", 1)))

    return shape


def render_line(slide, layout: Layout, element: Element):
    left, top, width, height = element_to_inches(layout, element)
    connector = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, left, top, left + width, top + height)
    connector.name = element.id
    connector.line.color.rgb = hex_to_rgb(style_value(element.style, "stroke", "#000000"))
    connector.line.width = Pt(float(style_value(element.style, "stroke_width", 1)))
    return connector
