from __future__ import annotations

from pptx.enum.text import PP_ALIGN, MSO_AUTO_SIZE, MSO_VERTICAL_ANCHOR
from pptx.util import Pt

from src.ir.normalize import element_to_inches, px_to_inches
from src.ir.schema import Element, Layout
from src.renderers.utils import hex_to_rgb, style_value


ALIGNMENTS = {
    "left": PP_ALIGN.LEFT,
    "center": PP_ALIGN.CENTER,
    "right": PP_ALIGN.RIGHT,
    "justify": PP_ALIGN.JUSTIFY,
}

VERTICAL_ANCHORS = {
    "top": MSO_VERTICAL_ANCHOR.TOP,
    "middle": MSO_VERTICAL_ANCHOR.MIDDLE,
    "bottom": MSO_VERTICAL_ANCHOR.BOTTOM,
}


def render_text(slide, layout: Layout, element: Element):
    left, top, width, height = element_to_inches(layout, element)
    textbox = slide.shapes.add_textbox(left, top, width, height)
    textbox.name = element.id
    text_frame = textbox.text_frame
    text_frame.clear()
    text_frame.word_wrap = bool(style_value(element.style, "word_wrap", True))
    text_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    text_frame.vertical_anchor = VERTICAL_ANCHORS.get(
        style_value(element.style, "vertical_anchor", "top"),
        MSO_VERTICAL_ANCHOR.TOP,
    )
    margin = Pt(float(style_value(element.style, "margin_pt", 0)))
    text_frame.margin_left = margin
    text_frame.margin_right = margin
    text_frame.margin_top = margin
    text_frame.margin_bottom = margin

    paragraph = text_frame.paragraphs[0]
    paragraph.alignment = ALIGNMENTS.get(style_value(element.style, "align", "left"), PP_ALIGN.LEFT)
    paragraph.line_spacing = float(style_value(element.style, "line_spacing", 1.0))

    runs = element.runs or []
    if not runs:
        run = paragraph.add_run()
        run.text = element.text or ""
        apply_run_style(run, layout, element.style)
    else:
        for source_run in runs:
            run = paragraph.add_run()
            run.text = source_run.text
            merged = {**element.style}
            if source_run.bold is not None:
                merged["bold"] = source_run.bold
            if source_run.italic is not None:
                merged["italic"] = source_run.italic
            if source_run.font_size is not None:
                merged["font_size"] = source_run.font_size
            if source_run.color is not None:
                merged["color"] = source_run.color
            apply_run_style(run, layout, merged)

    return textbox


def apply_run_style(run, layout: Layout, style: dict) -> None:
    font = run.font
    font.name = style_value(style, "font_family", "Aptos")
    font.size = Pt(font_size_to_points(layout, float(style_value(style, "font_size", 18))))
    font.bold = bool(style_value(style, "bold", False))
    font.italic = bool(style_value(style, "italic", False))
    font.color.rgb = hex_to_rgb(style_value(style, "color", "#222222"))


def font_size_to_points(layout: Layout, font_size_px: float) -> float:
    return px_to_inches(font_size_px, layout.canvas.height, layout.slide_size.height_in) * 72
