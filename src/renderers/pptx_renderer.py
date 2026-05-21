from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.util import Inches

from src.ir.normalize import sorted_elements
from src.ir.schema import Element, ElementType, Layout
from src.renderers.image_renderer import render_image
from src.renderers.shape_renderer import render_line, render_shape
from src.renderers.text_renderer import render_text


class PptxRenderer:
    def __init__(self, layout: Layout, base_dir: str | Path | None = None):
        self.layout = layout
        self.base_dir = Path(base_dir or ".").resolve()

    def render(self, output_path: str | Path) -> Path:
        output = Path(output_path)
        output.parent.mkdir(parents=True, exist_ok=True)

        prs = Presentation()
        prs.slide_width = Inches(self.layout.slide_size.width_in)
        prs.slide_height = Inches(self.layout.slide_size.height_in)
        blank_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(blank_layout)

        for element in sorted_elements(self.layout.elements):
            self._render_element(slide, element)

        prs.save(output)
        return output

    def _render_element(self, slide, element: Element) -> None:
        if element.type == ElementType.GROUP:
            for child in sorted_elements(element.children):
                self._render_element(slide, child)
            return
        if element.type == ElementType.TEXT:
            render_text(slide, self.layout, element)
            return
        if element.type == ElementType.SHAPE:
            render_shape(slide, self.layout, element)
            return
        if element.type == ElementType.LINE:
            render_line(slide, self.layout, element)
            return
        if element.type == ElementType.IMAGE:
            render_image(slide, self.layout, element, self.base_dir)
            return
        raise ValueError(f"unsupported element type: {element.type}")
